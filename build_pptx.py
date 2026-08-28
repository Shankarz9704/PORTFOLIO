import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Theme colors
    DARK_BG = RGBColor(15, 23, 42)       # #0f172a
    ACCENT_GREEN = RGBColor(5, 150, 105)  # #059669
    ACCENT_BLUE = RGBColor(37, 99, 235)   # #2563eb
    TEXT_LIGHT = RGBColor(248, 250, 252)  # #f8fafc
    TEXT_MUTED = RGBColor(148, 163, 184)  # #94a3b8
    CARD_BG = RGBColor(30, 41, 59)       # #1e293b
    WHITE = RGBColor(255, 255, 255)
    LINK_COLOR = RGBColor(56, 189, 248)  # #38bdf8

    def add_background(slide, color=DARK_BG):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = cat_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        p.font.name = "Arial"

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_LIGHT
        p2.font.name = "Arial"

    def add_clickable_link(paragraph, text, url, font_size=12, color=LINK_COLOR, bold=True):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.underline = True
        run.hyperlink.address = url
        return run

    # ================= SLIDE 1: Title Slide =================
    slide1 = prs.slides.add_slide(blank_layout)
    add_background(slide1, DARK_BG)

    # Accent decorative bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # Main Title
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.6), Inches(11), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PORTFOLIO WEBSITE PRESENTATION"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TEXT_LIGHT
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Bapanapalli Shankar — Full Stack & AI Data Analyst"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.name = "Arial"
    p2.space_before = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = "Funngro Portfolio Website Creation Project Submission"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Arial"
    p3.space_before = Pt(10)

    # Badge 1: Live Portfolio Website
    badge1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.8), Inches(10.8), Inches(0.75))
    badge1.fill.solid()
    badge1.fill.fore_color.rgb = CARD_BG
    badge1.line.color.rgb = ACCENT_GREEN
    tf_b1 = badge1.text_frame
    tf_b1.word_wrap = True
    p_b1 = tf_b1.paragraphs[0]
    p_b1.alignment = PP_ALIGN.LEFT
    r1 = p_b1.add_run()
    r1.text = " 🌐  Live Website:  "
    r1.font.size = Pt(13)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    add_clickable_link(p_b1, "https://bapanapallishankarportfolio.netlify.app/", "https://bapanapallishankarportfolio.netlify.app/", font_size=13)

    # Badge 2: Full Screen Web Presentation Link
    badge2 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.75), Inches(10.8), Inches(0.75))
    badge2.fill.solid()
    badge2.fill.fore_color.rgb = CARD_BG
    badge2.line.color.rgb = ACCENT_BLUE
    tf_b2 = badge2.text_frame
    tf_b2.word_wrap = True
    p_b2 = tf_b2.paragraphs[0]
    p_b2.alignment = PP_ALIGN.LEFT
    r2 = p_b2.add_run()
    r2.text = " 🖥️  Full-Screen Web Presentation:  "
    r2.font.size = Pt(13)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    add_clickable_link(p_b2, "https://bapanapallishankarportfolio.netlify.app/presentation.html", "https://bapanapallishankarportfolio.netlify.app/presentation.html", font_size=13)

    # ================= SLIDE 2: Executive Summary =================
    slide2 = prs.slides.add_slide(blank_layout)
    add_background(slide2, DARK_BG)
    add_header(slide2, "Executive Overview", "Project Summary & Clickable Live Links")

    # Card 1: Clickable Live Links
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_BLUE
    tf1 = card1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "🌐 Direct Clickable Links"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    links_list = [
        ("🖥️ Full-Screen Web PPT", "https://bapanapallishankarportfolio.netlify.app/presentation.html"),
        ("🚀 Netlify Live Website", "https://bapanapallishankarportfolio.netlify.app/"),
        ("🌐 GitHub Pages Site", "https://shankarz9704.github.io/PORTFOLIO/"),
        ("🐙 GitHub Repository", "https://github.com/Shankarz9704/PORTFOLIO"),
        ("📄 Submission PDF Document", "https://github.com/Shankarz9704/PORTFOLIO/blob/main/Portfolio_Website_Submission_Bapanapalli_Shankar.pdf")
    ]

    for label, url in links_list:
        p_l = tf1.add_paragraph()
        p_l.space_before = Pt(8)
        r_lbl = p_l.add_run()
        r_lbl.text = f"• {label}:\n   "
        r_lbl.font.size = Pt(11)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = TEXT_LIGHT
        
        add_clickable_link(p_l, url, url, font_size=10)

    # Card 2: Objectives Achieved
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = ACCENT_GREEN
    tf2 = card2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "🎯 Key Deliverables Achieved"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    achievements = [
        "Fully Responsive Web Design across Mobile, Tablet, and Desktop",
        "Interactive Dark/Light Theme Switching with persistent state",
        "Dynamic Animated Typewriter displaying core developer roles",
        "Project Showcase Gallery with category filter tabs",
        "Interactive Contact Form with instant toast feedback",
        "Production Build deployed on Netlify & GitHub Pages"
    ]
    for ach in achievements:
        p_a = tf2.add_paragraph()
        p_a.text = f"✓  {ach}"
        p_a.font.size = Pt(12)
        p_a.font.color.rgb = TEXT_LIGHT
        p_a.space_before = Pt(8)

    # ================= SLIDE 3: Developer Profile =================
    slide3 = prs.slides.add_slide(blank_layout)
    add_background(slide3, DARK_BG)
    add_header(slide3, "Candidate Profile", "Bapanapalli Shankar — Skills & Background")

    col_width = Inches(3.64)
    col_gap = Inches(0.4)

    # Col 1: Education
    c1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), col_width, Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.fill.background()
    t1 = c1.text_frame
    t1.word_wrap = True
    p = t1.paragraphs[0]
    p.text = "🎓 Education & Target"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    ed_text = [
        "B.Tech in Computer Science & Engineering",
        "Aggregate Marks: 75%",
        "Focus Areas: Software Engineering, Database Systems, Web Technologies, AI Analytics",
        "Target Roles: Python Developer, Java Full Stack, Spring Boot Engineer, Data Analyst"
    ]
    for t in ed_text:
        p_sub = t1.add_paragraph()
        p_sub.text = f"• {t}"
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_LIGHT
        p_sub.space_before = Pt(8)

    # Col 2: Technical Skills
    c2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + col_width + col_gap, Inches(1.8), col_width, Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.fill.background()
    t2 = c2.text_frame
    t2.word_wrap = True
    p = t2.paragraphs[0]
    p.text = "💻 Technical Stack"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    skills = [
        "Frontend: HTML5, CSS3, JavaScript (ES6+), CSS Grid/Flexbox",
        "Backend: Python, Java, Spring Boot, REST APIs",
        "Databases: MySQL, SQL Server, MongoDB",
        "Analytics: Data Processing, Business Intelligence, Python Pandas",
        "Tools: Git, GitHub, Netlify, VS Code"
    ]
    for s in skills:
        p_sub = t2.add_paragraph()
        p_sub.text = f"• {s}"
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_LIGHT
        p_sub.space_before = Pt(8)

    # Col 3: Experience
    c3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + (col_width + col_gap)*2, Inches(1.8), col_width, Inches(5.0))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.fill.background()
    t3 = c3.text_frame
    t3.word_wrap = True
    p = t3.paragraphs[0]
    p.text = "🏢 Work Experience"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    exps = [
        "Concentrix — Operations & Service Delivery Specialist",
        "Altruist Customer Management — Customer Support Associate",
        "Rooman Technologies — Technical Training & Development",
        "TaPTaP — Technical Support & System Handling",
        "ExcelR — Data Science & Analytics Training"
    ]
    for e in exps:
        p_sub = t3.add_paragraph()
        p_sub.text = f"• {e}"
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_LIGHT
        p_sub.space_before = Pt(8)

    # ================= SLIDE 4: Architecture & Features =================
    slide4 = prs.slides.add_slide(blank_layout)
    add_background(slide4, DARK_BG)
    add_header(slide4, "Technical Architecture", "Website Structure & Design System")

    w_card = Inches(5.6)
    h_card = Inches(2.3)

    features = [
        ("📱 Fluid Responsive Layout", "Mobile-first responsive architecture using pure CSS Grid and Flexbox. Adapts seamlessly across smartphones, tablets, laptops, and ultra-wide displays."),
        ("🌙 Dark & Light Theme System", "Toggleable visual themes with instant CSS custom property switching and persistent setting saved in browser localStorage."),
        ("⚡ Dynamic Typewriter & Animation", "Custom JavaScript engine driving an animated text hero switcher cycling developer specializations without external library bloat."),
        ("📊 Filterable Projects Showcase", "Interactive category tabs filtering projects dynamically by Full Stack, AI Data Analytics, Web Apps, and Security.")
    ]

    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.4)),
        (Inches(6.8), Inches(4.4))
    ]

    for idx, (title, desc) in enumerate(features):
        x, y = coords[idx]
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w_card, h_card)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_BLUE if idx % 2 == 0 else ACCENT_GREEN
        tf_f = card.text_frame
        tf_f.word_wrap = True
        
        p = tf_f.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN if idx % 2 == 0 else ACCENT_BLUE

        p_d = tf_f.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11.5)
        p_d.font.color.rgb = TEXT_LIGHT
        p_d.space_before = Pt(4)

    # ================= SLIDE 5: Featured Projects =================
    slide5 = prs.slides.add_slide(blank_layout)
    add_background(slide5, DARK_BG)
    add_header(slide5, "Projects Showcase", "Key Software & Analytics Projects")

    p_cards = [
        ("📲 Mobile Recharge Application", "Full Stack Web Application", "Java, Spring Boot, MySQL, REST API", "Built an enterprise-grade online mobile recharge system featuring user authentication, plan management, and instant payment transaction status processing.", "https://github.com/Shankarz9704"),
        ("📈 AI Data Analytics Dashboard", "Business Intelligence & Data Science", "Python, Pandas, SQL, BI Visuals", "Developed analytical pipelines and business intelligence dashboards to process structured datasets, perform statistical modeling, and generate actionable insights.", "https://github.com/Shankarz9704"),
        ("🌐 Personal Portfolio Website", "Personal Brand & Funngro Submission", "HTML5, CSS3, ES6 JS, Netlify", "Designed and deployed a modern interactive portfolio platform featuring theme switching, project filters, submission documentation, and live hosting.", "https://bapanapallishankarportfolio.netlify.app/")
    ]

    p_width = Inches(3.64)
    p_gap = Inches(0.4)

    for idx, (p_title, p_sub, p_tech, p_desc, p_url) in enumerate(p_cards):
        x = Inches(0.8) + idx * (p_width + p_gap)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), p_width, Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GREEN if idx == 2 else ACCENT_BLUE
        tf_p = card.text_frame
        tf_p.word_wrap = True

        p = tf_p.paragraphs[0]
        p.text = p_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p_sub_t = tf_p.add_paragraph()
        p_sub_t.text = p_sub
        p_sub_t.font.size = Pt(11.5)
        p_sub_t.font.bold = True
        p_sub_t.font.color.rgb = ACCENT_GREEN
        p_sub_t.space_before = Pt(4)

        p_tech_t = tf_p.add_paragraph()
        p_tech_t.text = f"Stack: {p_tech}"
        p_tech_t.font.size = Pt(10.5)
        p_tech_t.font.color.rgb = ACCENT_BLUE
        p_tech_t.space_before = Pt(4)

        p_desc_t = tf_p.add_paragraph()
        p_desc_t.text = p_desc
        p_desc_t.font.size = Pt(11.5)
        p_desc_t.font.color.rgb = TEXT_LIGHT
        p_desc_t.space_before = Pt(8)

        p_link_t = tf_p.add_paragraph()
        p_link_t.space_before = Pt(10)
        r_link_lbl = p_link_t.add_run()
        r_link_lbl.text = "🔗 Link: "
        r_link_lbl.font.size = Pt(10.5)
        r_link_lbl.font.bold = True
        r_link_lbl.font.color.rgb = WHITE
        add_clickable_link(p_link_t, "Access Project", p_url, font_size=10.5)

    # ================= SLIDE 6: Conclusion & Contact =================
    slide6 = prs.slides.add_slide(blank_layout)
    add_background(slide6, DARK_BG)
    add_header(slide6, "Contact & Easy Access Links", "Developer Contact & Complete Link Directory")

    card_c = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.33), Inches(5.0))
    card_c.fill.solid()
    card_c.fill.fore_color.rgb = CARD_BG
    card_c.line.color.rgb = ACCENT_GREEN
    tf_c = card_c.text_frame
    tf_c.word_wrap = True

    p = tf_c.paragraphs[0]
    p.text = "👨‍💻 Bapanapalli Shankar"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p_sub = tf_c.add_paragraph()
    p_sub.text = "Python & Java Full Stack Developer | AI Data Analyst"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = ACCENT_GREEN
    p_sub.space_before = Pt(2)

    contacts = [
        ("🖥️ Full-Screen Web Presentation", "https://bapanapallishankarportfolio.netlify.app/presentation.html"),
        ("🌐 Live Portfolio Website", "https://bapanapallishankarportfolio.netlify.app/"),
        ("🐙 GitHub Repository", "https://github.com/Shankarz9704/PORTFOLIO"),
        ("💼 LinkedIn Profile", "https://www.linkedin.com/in/shankar9704"),
        ("📧 Email Contact", "mailto:shankarcode291439@gmail.com"),
        ("📞 Phone", "+91 7842161185")
    ]

    for label, val in contacts:
        p_item = tf_c.add_paragraph()
        p_item.space_before = Pt(8)
        
        r_l = p_item.add_run()
        r_l.text = f"{label}:  "
        r_l.font.size = Pt(12)
        r_l.font.bold = True
        r_l.font.color.rgb = TEXT_LIGHT

        if val.startswith("http") or val.startswith("mailto"):
            add_clickable_link(p_item, val, val, font_size=12)
        else:
            r_v = p_item.add_run()
            r_v.text = val
            r_v.font.size = Pt(12)
            r_v.font.color.rgb = LINK_COLOR

    output_path = r"d:\MOBILE RECHARGE APPLICATION\funngro\Portfolio_Presentation_Bapanapalli_Shankar.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
